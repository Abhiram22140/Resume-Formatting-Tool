from pptx import Presentation
from pptx.dml.color import RGBColor

def _copy_font_properties(src_run, dest_run):
    src_font = src_run.font
    dst_font = dest_run.font

    if src_font.size:
        dst_font.size = src_font.size
    if src_font.name:
        dst_font.name = src_font.name

    dst_font.bold = src_font.bold
    dst_font.italic = src_font.italic

    if src_font.color:
        try:
            rgb = src_font.color.rgb
            dst_font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
        except AttributeError:
            pass

def _replace_text_preserve_format(shape, new_text: str):
    if not shape.has_text_frame:
        return

    tf = shape.text_frame

    src_run = None
    if tf.paragraphs and tf.paragraphs[0].runs:
        src_run = tf.paragraphs[0].runs[0]

    tf.clear()

    for i, line in enumerate(new_text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if src_run:
            _copy_font_properties(src_run, run)

def merge_into_template(parsed: dict, template_pptx_path: str, output_path: str):
    prs = Presentation(template_pptx_path)
    slide = prs.slides[0]

    # Prepare experience text
    exp_lines = []
    for exp in parsed.get("experience", []):
        header = f"{exp.get('position', '')}, {exp.get('company', '')} ({exp.get('dates', '')})".strip()
        if exp.get("description"):
            exp_lines.append(header)
            for line in exp["description"].split("\n"):
                exp_lines.append(f"• {line.strip('• ').strip()}")
            exp_lines.append("")
        else:
            exp_lines.append(header)
            exp_lines.append("")
    experience_text = "\n".join(exp_lines).strip()

    # Prepare education text
    edu_lines = []
    for edu in parsed.get("education", []):
        line = f"{edu.get('degree', '')}, {edu.get('institution', '')} ({edu.get('start', '')} – {edu.get('end', '')})"
        edu_lines.append(line)
    education_text = "\n".join(edu_lines).strip()

    content_map = {
        "Name": parsed.get("name", ""),
        "Role": parsed.get("role", ""),
        "Email": parsed.get("email", ""),
        "Phone": parsed.get("phone", ""),
        "Address": parsed.get("address", ""),
        "Skills": ", ".join(parsed.get("skills", [])),
        "Summary": parsed.get("summary", ""),
        "Experience": experience_text,
        "Education": education_text
    }

    for shape in slide.shapes:
        if shape.has_text_frame:
            current_text = shape.text_frame.text.strip()
            for key, value in content_map.items():
                if current_text.lower() == key.lower():
                    _replace_text_preserve_format(shape, value)
                    break  # Stop after first match

    prs.save(output_path)
